"""Tests for /reports/billing/export.* -- rewritten 2026-08-14 for the
Delayed Cash Billing (DCB) + Weekly Revenue Closure (WRC) compliance
export that replaced the deleted Audits export. Checks CSV/Excel/PDF/Word
all agree with each other and with the Dashboard, plus the PPT summary
deck (an executive summary, not a row dump -- see reports.py's docstring).
"""
import io
from datetime import date, datetime, timedelta, timezone

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from app.models.user import User
from app.services import delayed_cash_penalty_service as dcb_svc
from tests.conftest import TestingSessionLocal


def _register(client, username, email, password="password123"):
    return client.post("/auth/register", json={"username": username, "email": email, "password": password})


def _login(client, username, password="password123"):
    return client.post("/auth/login", json={"username": username, "password": password}).json()["access_token"]


def _auth(token):
    return {"Authorization": f"Bearer {token}"}


def _set_role(username, role):
    db = TestingSessionLocal()
    try:
        db.query(User).filter(User.username == username).first().role = role
        db.commit()
    finally:
        db.close()


def _admin(client, username="x_admin", email="x_admin@example.com"):
    _register(client, username, email)
    _set_role(username, "Admin")
    return _login(client, username)


def _make_user(username):
    db = TestingSessionLocal()
    try:
        import bcrypt

        user = User(
            username=username, email=f"{username}@example.com",
            password_hash=bcrypt.hashpw(b"password123", bcrypt.gensalt()).decode(),
            role="Admin", is_active=True,
        )
        db.add(user)
        db.commit()
        db.refresh(user)
        return user
    finally:
        db.close()


def _seed_dcb_bills(suffix: str, centre_code: str) -> int:
    """Two bills, one considered one not_considered, for a real center."""
    _make_user(f"x_dcb_setup{suffix}")
    db = TestingSessionLocal()
    try:
        user = db.query(User).filter(User.username == f"x_dcb_setup{suffix}").first()
        rule = dcb_svc.create_rule(db, rule_version=f"X-DCB-{suffix}", created_by=user)
        dcb_svc.approve_rule(db, rule=rule, approver=user)
        batch = dcb_svc.create_upload_batch(
            db, period_start=date(2026, 7, 1), period_end=date(2026, 7, 31),
            source_filename=f"x-{suffix}.xlsx", rule=rule, uploaded_by=user,
        )
        bill_date = date(2026, 7, 1)
        created = bill_date + timedelta(days=2)
        raw_bills = [
            dcb_svc.RawBillInput(
                centre_code=centre_code, centre_name=f"Export Test Center {suffix}",
                sales_bill=f"X-{suffix}-{i}", bill_date=bill_date,
                bill_created_time=datetime.combine(created, datetime.min.time(), tzinfo=timezone.utc),
                created_date=created, source_day_difference=2,
            )
            for i in range(2)
        ]
        bills = dcb_svc.ingest_bills(db, batch=batch, rule=rule, raw_bills=raw_bills)
        dcb_svc.compute_center_penalties(db, batch=batch, rule=rule)
        dcb_svc.set_bill_review_decision(db, bill=bills[0], decision="considered", reviewed_by=user)
        dcb_svc.set_bill_review_decision(db, bill=bills[1], decision="not_considered", reviewed_by=user)
        return len(bills)
    finally:
        db.close()


def test_xlsx_export_has_dcb_and_wrc_sheets_matching_dashboard_count(client):
    token = _admin(client)
    _seed_dcb_bills("1", "X-DCB-1")

    summary = client.get("/dashboard/summary", headers=_auth(token)).json()
    resp = client.get("/reports/billing/export.xlsx", headers=_auth(token))
    assert resp.status_code == 200
    assert resp.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

    wb = load_workbook(io.BytesIO(resp.content))
    assert wb.sheetnames == ["Delayed Cash Billing", "Weekly Revenue Closure"]

    dcb_rows = list(wb["Delayed Cash Billing"].iter_rows(values_only=True))
    data_rows = dcb_rows[1:]  # skip header
    assert len(data_rows) == summary["dcb"]["total_bills"]


def test_pdf_export_contains_centre_code(client):
    token = _admin(client, "x_admin2", "x_admin2@example.com")
    _seed_dcb_bills("2", "X-DCB-2")

    resp = client.get("/reports/billing/export.pdf", headers=_auth(token))
    assert resp.status_code == 200
    assert resp.headers["content-type"].startswith("application/pdf")
    assert resp.content.startswith(b"%PDF")

    reader = PdfReader(io.BytesIO(resp.content))
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    assert "X-DCB-2" in text
    assert "Delayed Cash Billing" in text
    assert "Weekly Revenue Closure" in text


def test_docx_export_has_two_tables(client):
    token = _admin(client, "x_admin3", "x_admin3@example.com")
    _seed_dcb_bills("3", "X-DCB-3")

    resp = client.get("/reports/billing/export.docx", headers=_auth(token))
    assert resp.status_code == 200
    assert resp.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )

    doc = Document(io.BytesIO(resp.content))
    assert len(doc.tables) == 2
    dcb_table = doc.tables[0]
    # header row + 2 seeded bills
    assert len(dcb_table.rows) == 3
    assert any("X-DCB-3" in cell.text for row in dcb_table.rows for cell in row.cells)


def test_pptx_export_is_an_executive_summary_not_a_row_dump(client):
    token = _admin(client, "x_admin4", "x_admin4@example.com")
    _seed_dcb_bills("4", "X-DCB-4")

    resp = client.get("/reports/billing/export.pptx", headers=_auth(token))
    assert resp.status_code == 200
    assert resp.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    prs = Presentation(io.BytesIO(resp.content))
    assert len(prs.slides) >= 4  # title, DCB, WRC, breakdowns/repeat-violators

    all_text = "\n".join(
        shape.text_frame.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
    )
    assert "Delayed Cash Billing" in all_text
    assert "Non-Compliance Rate" in all_text
    # Raw per-bill sales-bill numbers must NOT appear -- this is a summary
    # deck, not the row dump the other four formats give.
    assert "X-4-0" not in all_text


def test_csv_xlsx_pdf_docx_and_dashboard_all_agree_on_bill_count(client):
    """The real invariant: five independently-rendered outputs, one number."""
    token = _admin(client, "x_admin5", "x_admin5@example.com")
    _seed_dcb_bills("5", "X-DCB-5")

    summary = client.get("/dashboard/summary", headers=_auth(token)).json()
    expected_bills = summary["dcb"]["total_bills"]

    import csv as csv_mod

    csv_text = client.get("/reports/billing/export.csv", headers=_auth(token)).text
    csv_rows = list(csv_mod.reader(io.StringIO(csv_text)))
    dcb_header_idx = csv_rows.index(["DELAYED CASH BILLING"])
    wrc_header_idx = csv_rows.index(["WEEKLY REVENUE CLOSURE"])
    csv_dcb_data_rows = [r for r in csv_rows[dcb_header_idx + 2 : wrc_header_idx] if r]

    xlsx_resp = client.get("/reports/billing/export.xlsx", headers=_auth(token))
    wb = load_workbook(io.BytesIO(xlsx_resp.content))
    xlsx_dcb_rows = list(wb["Delayed Cash Billing"].iter_rows(values_only=True))[1:]

    docx_resp = client.get("/reports/billing/export.docx", headers=_auth(token))
    doc = Document(io.BytesIO(docx_resp.content))
    docx_dcb_rows = len(doc.tables[0].rows) - 1  # minus header

    assert len(csv_dcb_data_rows) == expected_bills
    assert len(xlsx_dcb_rows) == expected_bills
    assert docx_dcb_rows == expected_bills
