"""Billing compliance exports -- CSV/Excel/PDF/Word for the full raw row
set (Delayed Cash Billing bills + Weekly Revenue Closure incidents, each
with the decision, Vigilance's on-file remark, and whatever the center
itself submitted/was recorded as saying), and PPT as an executive-summary
deck built from the same compute_summary() the Dashboard shows -- a slide
full of a thousand raw rows is not a report, so PPT gets the KPIs and
breakdowns instead of the row dump the other four formats give.

Rewritten 2026-08-14 when the Audits/Findings domain this used to export
(app.models.audit) was deleted. Every number here still comes from
app.services.metrics -- this module never re-derives one.
"""
import io
from datetime import date

from fastapi import APIRouter, Depends, Response
from sqlalchemy.orm import Session

from app.auth import roles
from app.auth.dependencies import require_role
from app.database.database import get_db
from app.models.delayed_cash_billing import DelayedCashBill, DelayedCashCaseResponse, DelayedCashCenterPenalty
from app.models.user import User
from app.models.weekly_revenue_closure import WeeklyRevenueBillIncident
from app.services import report_service
from app.services import metrics
from app.services.metrics import MetricFilters, compute_summary

router = APIRouter(prefix="/reports", tags=["Reports"])

VIGILANCE_ROLES = (roles.ADMIN, roles.AUDITOR)

DCB_COLUMNS = [
    "Centre Code", "Centre Name", "Sales Bill", "Bill Date", "Decision",
    "Vigilance Remark", "Center Manager", "Manager NPID", "Center Remark", "Calculated Penalty",
]
WRC_COLUMNS = [
    "Centre Code", "Centre Name", "Incident Date", "Incident Type", "Decision",
    "Vigilance Remark", "Center Manager", "Manager NPID", "Center Remark",
]


def _latest_responses_by_center_penalty(db: Session, center_penalty_ids: list[int]) -> dict:
    if not center_penalty_ids:
        return {}
    rows = (
        db.query(DelayedCashCaseResponse)
        .filter(DelayedCashCaseResponse.center_penalty_id.in_(center_penalty_ids))
        .order_by(DelayedCashCaseResponse.center_penalty_id, DelayedCashCaseResponse.submitted_at.desc())
        .all()
    )
    latest: dict[int, DelayedCashCaseResponse] = {}
    for r in rows:
        latest.setdefault(r.center_penalty_id, r)  # first seen per id, since ordered desc
    return latest


def _center_penalty_ids_by_batch_and_centre(db: Session, bills: list[DelayedCashBill]) -> dict:
    """A bill has no FK to its case -- the link is (batch_id, centre_code)
    -- same resolution as _bills_with_center_penalty_ids in
    app/api/delayed_cash.py, inlined here since this needs the raw id,
    not a BillOut schema."""
    if not bills:
        return {}
    batch_ids = {b.batch_id for b in bills}
    centre_codes = {b.centre_code for b in bills}
    penalties = (
        db.query(DelayedCashCenterPenalty)
        .filter(DelayedCashCenterPenalty.batch_id.in_(batch_ids), DelayedCashCenterPenalty.centre_code.in_(centre_codes))
        .all()
    )
    return {(cp.batch_id, cp.centre_code): cp.id for cp in penalties}


def build_dcb_rows(db: Session, filters: MetricFilters) -> list[list]:
    """The ONE place that assembles a Delayed Cash Billing export's rows --
    every format below is a renderer over this same data."""
    bills = metrics.dcb_bills_query(db, filters).order_by(DelayedCashBill.id).all()
    cp_id_by_pair = _center_penalty_ids_by_batch_and_centre(db, bills)
    cp_ids = list(cp_id_by_pair.values())
    latest_by_cp = _latest_responses_by_center_penalty(db, cp_ids)

    rows = []
    for b in bills:
        cp_id = cp_id_by_pair.get((b.batch_id, b.centre_code))
        latest = latest_by_cp.get(cp_id) if cp_id else None
        rows.append(
            [
                b.centre_code,
                b.centre_name,
                b.sales_bill,
                b.bill_date.isoformat(),
                b.considered or "unreviewed",
                b.penalty_remarks or "",
                latest.responder_name if latest else "",
                latest.responder_npid if latest else "",
                latest.reason if latest else "",
                float(b.calculated_penalty),
            ]
        )
    return rows


def build_wrc_rows(db: Session, filters: MetricFilters) -> list[list]:
    incidents = metrics.wrc_incidents_query(db, filters).order_by(WeeklyRevenueBillIncident.id).all()
    return [
        [
            i.centre_code,
            i.centre_name,
            i.incident_date.isoformat(),
            i.mis_final_remark,
            i.considered or "unreviewed",
            i.penalty_remarks or "",
            i.center_manager or "",
            i.center_manager_npid or "",
            i.center_remarks or "",
        ]
        for i in incidents
    ]


def filters_from_query(period_from, period_to) -> MetricFilters:
    return MetricFilters(period_from=period_from, period_to=period_to)


def filters_to_dict(filters: MetricFilters) -> dict:
    return {
        "period_from": filters.period_from.isoformat() if filters.period_from else None,
        "period_to": filters.period_to.isoformat() if filters.period_to else None,
    }


def filters_from_dict(d: dict) -> MetricFilters:
    return MetricFilters(
        period_from=date.fromisoformat(d["period_from"]) if d.get("period_from") else None,
        period_to=date.fromisoformat(d["period_to"]) if d.get("period_to") else None,
    )


def render_csv(dcb_rows, wrc_rows) -> str:
    import csv

    buffer = io.StringIO()
    writer = csv.writer(buffer)
    writer.writerow(["DELAYED CASH BILLING"])
    writer.writerow(DCB_COLUMNS)
    writer.writerows(dcb_rows)
    writer.writerow([])
    writer.writerow(["WEEKLY REVENUE CLOSURE"])
    writer.writerow(WRC_COLUMNS)
    writer.writerows(wrc_rows)
    return buffer.getvalue()


def render_xlsx(dcb_rows, wrc_rows) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Delayed Cash Billing"
    ws1.append(DCB_COLUMNS)
    for row in dcb_rows:
        ws1.append(row)

    ws2 = wb.create_sheet("Weekly Revenue Closure")
    ws2.append(WRC_COLUMNS)
    for row in wrc_rows:
        ws2.append(row)

    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def _pdf_table(columns, rows):
    from reportlab.lib import colors
    from reportlab.platypus import Table, TableStyle

    table_data = [columns] + [[str(v) for v in row] for row in rows]
    table = Table(table_data, repeatRows=1)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f2937")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTSIZE", (0, 0), (-1, -1), 7),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ]
        )
    )
    return table


def render_pdf(dcb_rows, wrc_rows, title: str = "Billing Data Validation -- Compliance Export") -> bytes:
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=landscape(A4))
    styles = getSampleStyleSheet()

    elements = [
        Paragraph(title, styles["Title"]),
        Spacer(1, 12),
        Paragraph("Delayed Cash Billing", styles["Heading2"]),
        _pdf_table(DCB_COLUMNS, dcb_rows),
        PageBreak(),
        Paragraph("Weekly Revenue Closure", styles["Heading2"]),
        _pdf_table(WRC_COLUMNS, wrc_rows),
    ]
    doc.build(elements)
    return buffer.getvalue()


def render_docx(dcb_rows, wrc_rows) -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading("Billing Data Validation -- Compliance Export", level=0)

    doc.add_heading("Delayed Cash Billing", level=1)
    _docx_table(doc, DCB_COLUMNS, dcb_rows)

    doc.add_heading("Weekly Revenue Closure", level=1)
    _docx_table(doc, WRC_COLUMNS, wrc_rows)

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def _docx_table(doc, columns, rows):
    table = doc.add_table(rows=1, cols=len(columns))
    table.style = "Light Grid Accent 1"
    header_cells = table.rows[0].cells
    for idx, col in enumerate(columns):
        header_cells[idx].text = col
    for row in rows:
        cells = table.add_row().cells
        for idx, value in enumerate(row):
            cells[idx].text = str(value)


def render_pptx(db: Session, filters: MetricFilters) -> bytes:
    """An executive-summary deck, not a row dump -- built straight from
    compute_summary(), the same function the Dashboard renders."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    summary = compute_summary(db, filters)
    prs = Presentation()
    blank = prs.slide_layouts[6]

    def _slide():
        return prs.slides.add_slide(blank)

    def _title(slide, text):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        box.text_frame.text = text
        box.text_frame.paragraphs[0].font.size = Pt(28)
        box.text_frame.paragraphs[0].font.bold = True
        return box

    def _bullets(slide, lines, top=Inches(1.4)):
        box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(5))
        tf = box.text_frame
        for idx, line in enumerate(lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(18)

    s1 = _slide()
    _title(s1, "Billing Data Validation -- Compliance Summary")

    s2 = _slide()
    _title(s2, "Delayed Cash Billing")
    dcb = summary["dcb"]
    _bullets(
        s2,
        [
            f"Batches: {dcb['total_batches']}  ·  Bills: {dcb['total_bills']}",
            f"Considered: {dcb['considered']}  ·  Not Considered: {dcb['not_considered']}",
            f"Needs More Detail: {dcb['needs_more_detail']}  ·  Needs Proof: {dcb['needs_proof']}  ·  Unreviewed: {dcb['unreviewed']}",
            f"Non-Compliance Rate: {dcb['non_compliance_rate']}%" if dcb["non_compliance_rate"] is not None else "Non-Compliance Rate: not yet available",
            f"Total Validated Penalty: ₹{dcb['total_validated_penalty']:,.2f}",
        ],
    )

    s3 = _slide()
    _title(s3, "Weekly Revenue Closure")
    wrc = summary["wrc"]
    _bullets(
        s3,
        [
            f"Batches: {wrc['total_batches']}  ·  Incidents: {wrc['total_incidents']}",
            f"Considered: {wrc['considered']}  ·  Not Considered: {wrc['not_considered']}  ·  Unreviewed: {wrc['unreviewed']}",
            f"Non-Compliance Rate: {wrc['non_compliance_rate']}%" if wrc["non_compliance_rate"] is not None else "Non-Compliance Rate: not yet available",
            f"Total Center Penalty: ₹{wrc['total_center_penalty']:,.2f}  ·  Total Role Penalty: ₹{wrc['total_role_penalty']:,.2f}",
        ],
    )

    s4 = _slide()
    _title(s4, "Cluster / Zone Breakdown (Non-Compliant Centers)")
    cluster_lines = [f"{c['cluster']}: {c['non_compliant_center_count']} center(s)" for c in summary["cluster_breakdown"][:10]]
    zone_lines = [f"{z['zone']}: {z['non_compliant_center_count']} center(s)" for z in summary["zone_breakdown"][:10]]
    _bullets(s4, (cluster_lines or ["No cluster data yet"]) + [""] + (zone_lines or ["No zone data yet"]))

    s5 = _slide()
    _title(s5, "Repeat SOP Violators")
    repeat_lines = [f"{r['centre_name']} ({r['centre_code']}): {r['violation_count']} violations" for r in summary["repeated_centers"][:12]]
    if summary["repeated_centers_truncated"]:
        repeat_lines.append(f"(showing top {len(repeat_lines)} -- more exist)")
    _bullets(s5, repeat_lines or ["No repeat violators yet"])

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


MEDIA_TYPES = {
    "csv": "text/csv",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def export_response(db, user, filters: MetricFilters, fmt: str, *, name: str, template_id=None, regenerated_from_id=None):
    if fmt == "pptx":
        content = render_pptx(db, filters)
    else:
        dcb_rows = build_dcb_rows(db, filters)
        wrc_rows = build_wrc_rows(db, filters)
        if fmt == "csv":
            content = render_csv(dcb_rows, wrc_rows)
        elif fmt == "xlsx":
            content = render_xlsx(dcb_rows, wrc_rows)
        elif fmt == "pdf":
            content = render_pdf(dcb_rows, wrc_rows)
        elif fmt == "docx":
            content = render_docx(dcb_rows, wrc_rows)
        else:
            raise ValueError(f"Unknown export format {fmt!r}")

    report_service.record_history(
        db,
        name=name,
        template_id=template_id,
        filters_used=filters_to_dict(filters),
        format=fmt,
        generated_by=user,
        regenerated_from_id=regenerated_from_id,
    )

    return Response(
        content=content,
        media_type=MEDIA_TYPES[fmt],
        headers={"Content-Disposition": f"attachment; filename=billing_compliance_export.{fmt}"},
    )


@router.get("/billing/export.csv")
def export_billing_csv(
    period_from: date | None = None,
    period_to: date | None = None,
    db: Session = Depends(get_db),
    user: User = Depends(require_role(*VIGILANCE_ROLES)),
):
    return export_response(db, user, filters_from_query(period_from, period_to), "csv", name="Ad-hoc Billing Export")


@router.get("/billing/export.xlsx")
def export_billing_xlsx(
    period_from: date | None = None,
    period_to: date | None = None,
    db: Session = Depends(get_db),
    user: User = Depends(require_role(*VIGILANCE_ROLES)),
):
    return export_response(db, user, filters_from_query(period_from, period_to), "xlsx", name="Ad-hoc Billing Export")


@router.get("/billing/export.pdf")
def export_billing_pdf(
    period_from: date | None = None,
    period_to: date | None = None,
    db: Session = Depends(get_db),
    user: User = Depends(require_role(*VIGILANCE_ROLES)),
):
    return export_response(db, user, filters_from_query(period_from, period_to), "pdf", name="Ad-hoc Billing Export")


@router.get("/billing/export.docx")
def export_billing_docx(
    period_from: date | None = None,
    period_to: date | None = None,
    db: Session = Depends(get_db),
    user: User = Depends(require_role(*VIGILANCE_ROLES)),
):
    return export_response(db, user, filters_from_query(period_from, period_to), "docx", name="Ad-hoc Billing Export")


@router.get("/billing/export.pptx")
def export_billing_pptx(
    period_from: date | None = None,
    period_to: date | None = None,
    db: Session = Depends(get_db),
    user: User = Depends(require_role(*VIGILANCE_ROLES)),
):
    return export_response(db, user, filters_from_query(period_from, period_to), "pptx", name="Ad-hoc Billing Export")
